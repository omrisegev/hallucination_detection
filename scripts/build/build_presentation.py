"""
Build advisor meeting PPTX -- May 18, 2026
Spectral Hallucination Detection: Progress Update
Revised: removed multiview slides, added domain examples, updated competitor tables
"""
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.patches as mpatches
from matplotlib.lines import Line2D
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# -- paths
BASE  = os.path.dirname(os.path.abspath(__file__))
PLOTS = os.path.join(BASE, "presentation_plots")
OUT   = os.path.join(BASE, "Hallucination_Detection_May18.pptx")
TMP   = os.path.join(PLOTS, "_tmp")
os.makedirs(TMP, exist_ok=True)

# -- colors
NAVY   = RGBColor(0x1F, 0x49, 0x7D)
DGREY  = RGBColor(0x32, 0x32, 0x32)
BLUE   = RGBColor(0x00, 0x70, 0xC0)
ORANGE = RGBColor(0xFF, 0x66, 0x00)
GREEN  = RGBColor(0x00, 0x96, 0x40)
RED    = RGBColor(0xC0, 0x00, 0x00)
LGREY  = RGBColor(0xF2, 0xF2, 0xF2)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
PURPLE = RGBColor(0x70, 0x30, 0xA0)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
blank = prs.slide_layouts[6]

# =======================================================================
# HELPERS
# =======================================================================

def add_slide():
    return prs.slides.add_slide(blank)

def txb(slide, text, l, t, w, h, size=18, bold=False, color=DGREY,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    return tb

def title_bar(slide, text, subtitle=None):
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.33), Inches(1.15))
    bar.fill.solid()
    bar.fill.fore_color.rgb = NAVY
    bar.line.fill.background()
    txb(slide, text, 0.25, 0.08, 12.5, 0.9,
        size=28, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        txb(slide, subtitle, 0.25, 0.82, 12.5, 0.35,
            size=13, color=RGBColor(0xCC, 0xD8, 0xF0), align=PP_ALIGN.LEFT)

def img(slide, path, l, t, w, h=None):
    if h is None:
        return slide.shapes.add_picture(path, Inches(l), Inches(t), width=Inches(w))
    return slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w), Inches(h))

def bullet_box(slide, items, l, t, w, h, size=16):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    for i, (marker, text, color) in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"{marker}  {text}"
        run.font.size  = Pt(size)
        run.font.color.rgb = color

def simple_table(slide, headers, rows, l, t, w, h,
                 hdr_bg=NAVY, hdr_fg=WHITE,
                 row_bg=LGREY, row_fg=DGREY, alt_bg=WHITE,
                 col_widths=None, hdr_size=12, row_size=11):
    n_cols = len(headers)
    n_rows = len(rows)
    tbl = slide.shapes.add_table(
        n_rows + 1, n_cols,
        Inches(l), Inches(t), Inches(w), Inches(h)).table
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
    for ci, hdr in enumerate(headers):
        cell = tbl.cell(0, ci)
        cell.fill.solid(); cell.fill.fore_color.rgb = hdr_bg
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = hdr
        run.font.size = Pt(hdr_size); run.font.bold = True
        run.font.color.rgb = hdr_fg
    for ri, row in enumerate(rows):
        bg = row_bg if ri % 2 == 0 else alt_bg
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            cell.fill.solid(); cell.fill.fore_color.rgb = bg
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT
            run = p.add_run(); run.text = str(val)
            run.font.size = Pt(row_size); run.font.color.rgb = row_fg
    return tbl

def accent_box(slide, text, l, t, w, h,
               bg=BLUE, fg=WHITE, size=14, bold=False):
    box = slide.shapes.add_shape(
        1, Inches(l), Inches(t), Inches(w), Inches(h))
    box.fill.solid(); box.fill.fore_color.rgb = bg
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(size); run.font.bold = bold
    run.font.color.rgb = fg

# =======================================================================
# CHART GENERATORS
# =======================================================================

def make_math_example():
    """Math problem structure + H(n) trajectory with feature annotations."""
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.8),
                             gridspec_kw={'width_ratios': [1, 1.5]})

    # LEFT: problem structure diagram
    ax = axes[0]
    ax.axis('off')
    ax.set_facecolor('#FAFAFA')
    for spine in ax.spines.values():
        spine.set_visible(True); spine.set_edgecolor('#CCCCCC')

    content = [
        (0.95, 'MATH-500: What Does H(n) Cover?', '#1F497D', 'bold', 11),
        (0.84, 'Q: "Solve x^2 + 3x + 2 = 0"', '#333', 'normal', 9),
        (0.74, '[Chain-of-Thought -- ~500 tokens]', '#1565C0', 'bold', 9),
        (0.66, '"I need to factor this quadratic..."', '#444', 'normal', 8.5),
        (0.58, '"(x+1)(x+2) = 0, so x = -1 or x = -2"', '#444', 'normal', 8.5),
        (0.50, '"Checking: (-1)^2+3(-1)+2=0 correct"', '#444', 'normal', 8.5),
        (0.40, '[Answer -- ~20 tokens]', '#2E7D32', 'bold', 9),
        (0.32, '"The solutions are x = -1 and x = -2."', '#2E7D32', 'normal', 8.5),
        (0.18, 'H(n) covers the FULL trace', '#777', 'italic', 8),
        (0.11, '(reasoning + answer, 500-2000 tokens total)', '#777', 'italic', 8),
    ]
    for y, txt, col, wt, sz in content:
        ax.text(0.05, y, txt, ha='left', va='top', fontsize=sz, color=col,
                fontweight=wt if wt != 'italic' else 'normal',
                fontstyle='italic' if wt == 'italic' else 'normal',
                transform=ax.transAxes)

    # RIGHT: simulated H(n) -- BOTH correct and hallucinating
    ax = axes[1]
    np.random.seed(42)
    n = 680
    cot_end = 580
    t = np.arange(n)

    # Correct trace: smooth reasoning, low stable answer
    base_c = 2.8 + 0.4 * np.sin(t * 0.019) + 0.3 * np.cos(t * 0.043)
    trace_c = base_c + np.random.randn(n) * 0.18
    trace_c[cot_end:] = 0.85 + np.random.randn(n - cot_end) * 0.12

    # Hallucinating trace: higher entropy, spikes, rising end (rpdi)
    np.random.seed(7)
    base_h = 3.5 + 0.6 * np.sin(t * 0.023) + 0.5 * np.cos(t * 0.037)
    trace_h = base_h + np.random.randn(n) * 0.35
    trace_h[88:96]   += 2.5
    trace_h[230:238] += 1.9
    trace_h[410:418] += 2.1
    trace_h[cot_end:] = 2.5 + np.random.randn(n - cot_end) * 0.28

    ax.plot(t, trace_c, color='#2E7D32', lw=1.0, alpha=0.88, label='Correct answer')
    ax.plot(t, trace_h, color='#C00000', lw=0.9, alpha=0.75, label='Hallucinated answer')
    ax.axvspan(0, cot_end, alpha=0.05, color='orange')
    ax.axvspan(cot_end, n, alpha=0.05, color='limegreen')
    ax.axvline(cot_end, color='grey', lw=1.5, ls='--', alpha=0.7)
    ax.text(cot_end * 0.5, 7.3, 'Chain-of-Thought', ha='center',
            fontsize=9, color='#E65C00', fontstyle='italic')
    ax.text(cot_end + (n - cot_end) * 0.5, 7.3, 'Answer', ha='center',
            fontsize=9, color='#2E7D32', fontstyle='italic')
    ax.annotate('cusum_max\n(regime shift)', xy=(91, trace_h[91]),
                xytext=(180, 6.3), fontsize=8.5, color='#C00000',
                arrowprops=dict(arrowstyle='->', color='#C00000', lw=1.3))
    ax.annotate('rpdi: end-entropy\nrises in halluc.',
                xy=(cot_end + 40, trace_h[cot_end + 40]),
                xytext=(cot_end - 110, 6.4), fontsize=8.5, color='#7030A0',
                arrowprops=dict(arrowstyle='->', color='#7030A0', lw=1.3))
    ax.set_xlabel('Token position n', fontsize=11)
    ax.set_ylabel('H(n) = -sum p log p', fontsize=11)
    ax.set_title('H(n): correct (green) vs hallucinated (red)\n(simulated; representative of MATH-500/Qwen-7B T=1.0 pattern)',
                 fontsize=10.5, fontweight='bold')
    ax.set_facecolor('#FAFAFA')
    ax.yaxis.grid(True, alpha=0.25)
    ax.legend(fontsize=10, loc='upper right')

    plt.tight_layout(pad=0.8)
    path = os.path.join(TMP, 'math_example.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def make_gpqa_bar():
    fig, ax = plt.subplots(figsize=(5, 3.5))
    models = ['7B models\n(Mistral/Qwen)', 'Qwen-72B AWQ']
    aucs   = [65.4, 69.0]
    bars = ax.bar(models, aucs, color=['#4472C4', '#FF6600'], width=0.45, zorder=3)
    ax.set_ylim(55, 75)
    ax.set_ylabel('AUROC (%)', fontsize=12)
    ax.set_title('GPQA Diamond: Scale Improves Detection',
                 fontsize=13, fontweight='bold')
    ax.axhline(70, color='green', lw=1.5, ls='--', zorder=2, label='70% target')
    ax.yaxis.grid(True, alpha=0.3, zorder=0); ax.set_axisbelow(True)
    for bar, v in zip(bars, aucs):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.3, f'{v:.1f}%',
                ha='center', va='bottom', fontsize=13, fontweight='bold')
    ax.legend(fontsize=10)
    ax.text(1, 57.5, '+3.6 pp (+5.5% relative)', ha='center',
            fontsize=10, color='#FF6600', fontstyle='italic')
    plt.tight_layout()
    path = os.path.join(TMP, 'gpqa_bar.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def make_rag_citation_example():
    fig, ax = plt.subplots(figsize=(11, 5.8))
    ax.axis('off')
    ax.set_facecolor('white'); fig.patch.set_facecolor('white')
    y = 0.97

    def line(text, dy, size=9.5, color='black', weight='normal', italic=False):
        nonlocal y
        ax.text(0.01, y, text, transform=ax.transAxes, fontsize=size,
                color=color, fontweight=weight,
                fontstyle='italic' if italic else 'normal', va='top')
        y -= dy

    line('HotpotQA -- How H(n) Is Sliced in RAG', 0.09, size=11,
         color='#1F497D', weight='bold')
    line('The model generates a response with citation markers [n].',
         0.06, size=9, color='#555', italic=True)
    line('We slice H(n) at each citation boundary => each statement gets its own spectral feature vector.',
         0.065, size=9, color='#555', italic=True)
    line('Q: "Were Coldplay and The Chainsmokers both founded in the same country?"',
         0.075, size=9.5, color='#333')
    line('Retrieved: [1] Coldplay = British rock band, London, 1996.  [2] Chainsmokers = American duo, NYC.',
         0.07, size=8.5, color='#555')
    line('Model output (with citation markers):',
         0.065, size=9, color='#555', weight='bold')

    # Statement 1 -- grounded
    ax.add_patch(mpatches.FancyBboxPatch(
        (0.01, y - 0.082), 0.98, 0.082,
        boxstyle='round,pad=0.01', facecolor='#E8F5E9', edgecolor='#2E7D32',
        lw=1.5, transform=ax.transAxes, zorder=2))
    ax.text(0.03, y - 0.012, '[GROUNDED]', transform=ax.transAxes,
            fontsize=7.5, color='#2E7D32', fontweight='bold', va='top')
    ax.text(0.16, y - 0.012,
            '"Coldplay is a British rock band formed in London. [1]"',
            transform=ax.transAxes, fontsize=9, color='#1B5E20', va='top')
    ax.text(0.16, y - 0.050,
            'H(n)[tok_start_1 : tok_end_1]  =>  spectral features extracted from this slice',
            transform=ax.transAxes, fontsize=8, color='#2E7D32', fontstyle='italic', va='top')
    y -= 0.10

    # Statement 2 -- grounded
    ax.add_patch(mpatches.FancyBboxPatch(
        (0.01, y - 0.082), 0.98, 0.082,
        boxstyle='round,pad=0.01', facecolor='#E8F5E9', edgecolor='#2E7D32',
        lw=1.5, transform=ax.transAxes, zorder=2))
    ax.text(0.03, y - 0.012, '[GROUNDED]', transform=ax.transAxes,
            fontsize=7.5, color='#2E7D32', fontweight='bold', va='top')
    ax.text(0.16, y - 0.012,
            '"The Chainsmokers are an American DJ duo from New York. [2]"',
            transform=ax.transAxes, fontsize=9, color='#1B5E20', va='top')
    ax.text(0.16, y - 0.050,
            'H(n)[tok_start_2 : tok_end_2]  =>  spectral features extracted from this slice',
            transform=ax.transAxes, fontsize=8, color='#2E7D32', fontstyle='italic', va='top')
    y -= 0.10

    # Statement 3 -- ungrounded
    ax.add_patch(mpatches.FancyBboxPatch(
        (0.01, y - 0.095), 0.98, 0.095,
        boxstyle='round,pad=0.01', facecolor='#FFEBEE', edgecolor='#C62828',
        lw=1.5, transform=ax.transAxes, zorder=2))
    ax.text(0.03, y - 0.012, '[UNGROUNDED]', transform=ax.transAxes,
            fontsize=7.5, color='#C62828', fontweight='bold', va='top')
    ax.text(0.19, y - 0.012,
            '"Therefore, they were NOT founded in the same country. [1][2]"',
            transform=ax.transAxes, fontsize=9, color='#7F0000', va='top')
    ax.text(0.19, y - 0.048,
            'Fabricated conclusion -- not stated in either retrieved document',
            transform=ax.transAxes, fontsize=8, color='#C62828', fontstyle='italic', va='top')
    ax.text(0.19, y - 0.072,
            'H(n)[tok_start_3 : tok_end_3]  =>  spectral features extracted => higher rpdi, spectral_entropy',
            transform=ax.transAxes, fontsize=8, color='#7F0000', fontstyle='italic', va='top')
    y -= 0.115

    line('Label: GROUNDED if cited passage title appears in gold_supporting_facts; else UNGROUNDED',
         0.05, size=8.5, color='#555', italic=True)
    line('Each statement\'s spectral features are fused with Nadler independently.',
         0.05, size=8.5, color='#555', italic=True)

    plt.tight_layout(pad=0.3)
    path = os.path.join(TMP, 'rag_citation_example.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def make_agentic_example():
    """ReAct loop structure + per-step spectral scores + Phi_min aggregation."""
    fig, axes = plt.subplots(1, 2, figsize=(12, 4.8),
                             gridspec_kw={'width_ratios': [1.2, 1]})

    # LEFT: ReAct trace
    ax = axes[0]
    ax.axis('off')
    ax.set_facecolor('white')
    ax.text(0.5, 0.97, 'ReAct Agent Trace (Phase 11a)', ha='center', va='top',
            fontsize=11, fontweight='bold', color='#1F497D',
            transform=ax.transAxes)
    ax.text(0.5, 0.89,
            'Q: "Were Coldplay and The Chainsmokers founded in the same country?"',
            ha='center', va='top', fontsize=8.5, color='#333',
            fontstyle='italic', transform=ax.transAxes)

    steps_data = [
        ('Step 1', 'Thought: "I need to check where both bands are from..."',
         'Action: Search[Coldplay origin]',
         '#E3F2FD', '#1565C0'),
        ('Step 2', 'Thought: "Coldplay is from the UK. Now for Chainsmokers..."',
         'Action: Search[Chainsmokers origin]',
         '#E8F5E9', '#2E7D32'),
        ('Step 3', 'Thought: "Chainsmokers are from the USA. UK != USA, so No."',
         'Action: Finish[No, not the same country]',
         '#FFF3E0', '#E65C00'),
    ]

    y = 0.80
    for step_lbl, thought, action, bg, color in steps_data:
        rect = mpatches.FancyBboxPatch(
            (0.01, y - 0.21), 0.98, 0.21,
            boxstyle='round,pad=0.01', facecolor=bg, edgecolor=color, lw=1.5,
            transform=ax.transAxes, zorder=2)
        ax.add_patch(rect)
        ax.text(0.04, y - 0.03, step_lbl, ha='left', va='top', fontsize=8.5,
                fontweight='bold', color=color, transform=ax.transAxes)
        ax.text(0.19, y - 0.03, thought, ha='left', va='top', fontsize=7.8,
                color='#222', transform=ax.transAxes)
        ax.text(0.19, y - 0.13, action, ha='left', va='top', fontsize=7.8,
                color='#555', transform=ax.transAxes)
        y -= 0.26

    ax.text(0.5, 0.03,
            'H(n) extracted from Thought tokens only (50-150 tokens per step)',
            ha='center', va='bottom', fontsize=8, color='#555',
            fontstyle='italic', transform=ax.transAxes)

    # RIGHT: per-step Nadler scores + Phi_min
    ax = axes[1]
    scores = [0.74, 0.61, 0.79]
    phi_min = min(scores)
    bar_colors = ['#4472C4', '#C00000', '#4472C4']
    bars = ax.bar([1, 2, 3], scores, color=bar_colors, width=0.5, zorder=3)
    ax.axhline(phi_min, color='red', lw=2.0, ls='--', zorder=4,
               label=f'Phi_min = {phi_min:.2f} (weakest link)')
    ax.set_xticks([1, 2, 3])
    ax.set_xticklabels(['Step 1\nThought', 'Step 2\nThought', 'Step 3\nThought'], fontsize=10)
    ax.set_ylabel('Nadler spectral score', fontsize=11)
    ax.set_title('Spectral score per step\n=> Phi_min = min across steps',
                 fontsize=11, fontweight='bold')
    ax.set_ylim(0.45, 0.92)
    ax.yaxis.grid(True, alpha=0.3, zorder=0); ax.set_axisbelow(True)
    for bar, v in zip(bars, scores):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.01, f'{v:.2f}',
                ha='center', va='bottom', fontsize=12, fontweight='bold')
    ax.annotate('Uncertain step\n=> drives Phi_min',
                xy=(2, phi_min + 0.01), xytext=(2.45, 0.54),
                fontsize=9, color='red',
                arrowprops=dict(arrowstyle='->', color='red', lw=1.5))
    ax.legend(fontsize=9, loc='upper left')
    ax.set_facecolor('#FAFAFA')

    plt.tight_layout(pad=0.8)
    path = os.path.join(TMP, 'agentic_example.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def make_agentic_bar():
    fig, ax = plt.subplots(figsize=(5.5, 3.8))
    methods = ['AUQ Phi_min\n(Zhang 2026)\nALFWorld', 'Ours Phi_min\n(mid-run)\nDeepSeek/2Wiki']
    aucs    = [79.1, 85.0]
    bars = ax.bar(methods, aucs, color=['#4472C4', '#FF6600'], width=0.4, zorder=3)
    ax.set_ylim(70, 90)
    ax.set_ylabel('AUROC (%)', fontsize=12)
    ax.set_title('Agentic: Spectral Nadler vs AUQ\n(mid-run, not yet official)',
                 fontsize=11, fontweight='bold')
    ax.yaxis.grid(True, alpha=0.3, zorder=0); ax.set_axisbelow(True)
    for bar, v in zip(bars, aucs):
        ax.text(bar.get_x() + bar.get_width() / 2, v + 0.2, f'{v:.1f}%',
                ha='center', va='bottom', fontsize=14, fontweight='bold')
    ax.annotate('', xy=(1, 85.0), xytext=(1, 79.1),
                arrowprops=dict(arrowstyle='<->', color='green', lw=2))
    ax.text(1.22, 82, '+5.9 pp', color='green', fontsize=11, fontweight='bold')
    ax.text(0.5, 71.2,
            '* Mistral-24B + Qwen-72B cells still pending',
            ha='center', fontsize=8.5, color='grey', fontstyle='italic',
            transform=ax.get_xaxis_transform())
    plt.tight_layout()
    path = os.path.join(TMP, 'agentic_bar.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def make_feature_importance_grid():
    from PIL import Image, ImageDraw
    imgs = [
        ("Math-500",     "meta_analysis_cell07_out0.png"),
        ("GSM8K",        "meta_analysis_cell07_out1.png"),
        ("GPQA",         "meta_analysis_cell07_out2.png"),
        ("QA (factual)", "meta_analysis_cell07_out3.png"),
        ("RAG",          "meta_analysis_cell07_out4.png"),
    ]
    loaded  = [(lbl, Image.open(os.path.join(PLOTS, fn))) for lbl, fn in imgs]
    tw, th  = 480, 320
    resized = [(lbl, im.resize((tw, th), Image.LANCZOS)) for lbl, im in loaded]
    canvas  = Image.new('RGB', (3 * tw, 2 * th), 'white')
    for idx, (lbl, im) in enumerate(resized):
        r, c = divmod(idx, 3)
        canvas.paste(im, (c * tw, r * th))
    draw = ImageDraw.Draw(canvas)
    x6, y6 = 2 * tw, 1 * th
    draw.rectangle([x6, y6, x6 + tw, y6 + th], fill='#F0F4F8', outline='#CCCCCC')
    note_lines = [
        "Cross-domain verdict:",
        "",
        "cusum_max & sw_var_peak",
        "  => top-3 in ALL 5 domains",
        "",
        "epr dominates Math / GSM8K",
        "rpdi dominates QA / RAG",
        "",
        "pe_min: rank 17/17 everywhere",
        "  => REMOVED from feature set",
        "",
        "16 features total remain",
    ]
    for li, text in enumerate(note_lines):
        draw.text((x6 + 18, y6 + 26 + li * 22), text, fill='#1F497D')
    path = os.path.join(TMP, 'feature_importance_grid.png')
    canvas.save(path)
    return path


def make_results_overview():
    """All experiments -- unified bar chart with competitor diamond markers."""
    fig, ax = plt.subplots(figsize=(12, 5.4))

    tasks = [
        ('MATH-500 / Qwen-7B T=1.0',        90.0, '#4472C4', None),
        ('GSM8K / Llama-3.1-8B',             76.0, '#4472C4', 72.0),
        ('GPQA Diamond / Qwen-72B',          69.0, '#7030A0', None),
        ('TriviaQA CoT / Falcon-10B',        53.6, '#C00000', None),
        ('WebQ CoT / Falcon-10B',            61.9, '#C00000', None),
        ('RAG HotpotQA / Qwen-7B',           79.5, '#FF6600', 72.92),
        ('RAG 2Wiki / Qwen-7B',              80.5, '#FF6600', None),
        ('RAG HotpotQA / Llama-8B *',        87.7, '#FF3300', 72.92),
        ('Agentic 2Wiki / DeepSeek-7B **',   85.0, '#00963F', 79.1),
    ]

    y_pos  = list(range(len(tasks)))
    labels = [t[0] for t in tasks]
    aucs   = [t[1] for t in tasks]
    colors = [t[2] for t in tasks]

    bars = ax.barh(y_pos, aucs, color=colors, height=0.6, zorder=3)
    ax.set_yticks(y_pos); ax.set_yticklabels(labels, fontsize=9)
    ax.set_xlim(40, 97)
    ax.set_xlabel('AUROC (%)', fontsize=11)
    ax.set_title('All Experiments -- Spectral Nadler Fusion', fontsize=13, fontweight='bold')
    ax.xaxis.grid(True, alpha=0.25, zorder=0); ax.set_axisbelow(True)
    ax.axvline(70, color='#AAAAAA', lw=1.5, ls='--', zorder=2)
    ax.text(70.5, len(tasks) - 0.5, '70% gate', fontsize=8.5, color='#777')

    for i, (_, _, _, comp) in enumerate(tasks):
        if comp is not None:
            ax.plot(comp, i, 'kD', markersize=8, zorder=6)

    for bar, v in zip(bars, aucs):
        ax.text(v + 0.3, bar.get_y() + bar.get_height() / 2,
                f'{v:.1f}%', ha='left', va='center', fontsize=9, fontweight='bold')

    legend_els = [
        mpatches.Patch(color='#4472C4', label='Math reasoning'),
        mpatches.Patch(color='#7030A0', label='Science MCQ'),
        mpatches.Patch(color='#C00000', label='Factual QA (neg.)'),
        mpatches.Patch(color='#FF6600', label='RAG grounding'),
        mpatches.Patch(color='#FF3300', label='RAG best cell'),
        mpatches.Patch(color='#00963F', label='Agentic'),
        Line2D([0], [0], marker='D', color='k', lw=0, markersize=8, label='Competitor'),
    ]
    ax.legend(handles=legend_els, fontsize=8.5, loc='lower right')
    ax.text(0.0, -0.11,
            '* llama8b/HotpotQA best RAG cell   ** mid-run signal, not yet official',
            transform=ax.transAxes, fontsize=8, color='grey', fontstyle='italic')

    plt.tight_layout(pad=1.0)
    path = os.path.join(TMP, 'results_overview.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


def make_formula_panel():
    """4-box panel with key feature formulas."""
    fig, axes = plt.subplots(2, 2, figsize=(5.2, 3.6))
    fig.patch.set_facecolor('white')
    formulas = [
        ('cusum_max',
         r'max|cumsum($H(n) - \bar{H}$)|',
         'Cumulative sum of mean-centred\nentropy — detects regime shifts',
         '#1F497D'),
        ('sw_var_peak',
         r'$\max_i\,\mathrm{var}(H[i:i{+}w])$, $w{=}16$',
         'Peak local variance across\na sliding window',
         '#7030A0'),
        ('rpdi',
         r'$\frac{\mathrm{mean}(H[\text{last }20\%])}{\mathrm{mean}(H[\text{full}])}$',
         'End-of-trace entropy rise:\nhallucinations spike at end',
         '#FF6600'),
        ('spectral_entropy',
         r'$-\sum \tilde{p}(f)\log\tilde{p}(f)$',
         r'$\tilde{p}(f)$ = normalised PSD;' + '\nenergy spread across freq.',
         '#2E7D32'),
    ]
    for idx, (name, formula, desc, color) in enumerate(formulas):
        r, c = divmod(idx, 2)
        ax = axes[r, c]
        ax.axis('off')
        ax.add_patch(mpatches.FancyBboxPatch(
            (0.03, 0.03), 0.94, 0.94,
            boxstyle='round,pad=0.03', facecolor='#F4F7FB',
            edgecolor=color, lw=1.8, transform=ax.transAxes))
        ax.text(0.5, 0.87, name, ha='center', va='top', fontsize=9.5,
                fontweight='bold', color=color, transform=ax.transAxes)
        ax.text(0.5, 0.57, formula, ha='center', va='center', fontsize=8.5,
                color='#222', transform=ax.transAxes)
        ax.text(0.5, 0.20, desc, ha='center', va='bottom', fontsize=7.5,
                color='#555', transform=ax.transAxes, style='italic',
                multialignment='center')
    plt.tight_layout(pad=0.4)
    path = os.path.join(TMP, 'formula_panel.png')
    fig.savefig(path, dpi=150, bbox_inches='tight', facecolor='white')
    plt.close()
    return path


# =======================================================================
# PRE-GENERATE CHARTS
# =======================================================================
print("Generating helper charts...")
GPQA_BAR      = make_gpqa_bar()
MATH_EX       = make_math_example()
AGENTIC_EX    = make_agentic_example()
AGENTIC_BAR   = make_agentic_bar()
RAG_EXAMPLE   = make_rag_citation_example()
RESULTS_OVW   = make_results_overview()
FEAT_GRID     = make_feature_importance_grid()
FORMULA_PANEL = make_formula_panel()
print("Charts done.")

# =======================================================================
# SLIDES
# =======================================================================

# -----------------------------------------------------------------------
# SLIDE 1: TITLE
# -----------------------------------------------------------------------
sl = add_slide()
bg = sl.shapes.add_shape(1, Inches(0), Inches(0), W, H)
bg.fill.solid(); bg.fill.fore_color.rgb = NAVY; bg.line.fill.background()
txb(sl, "Spectral Hallucination Detection", 0.7, 1.5, 12.0, 1.2,
    size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txb(sl, "Progress Update -- May 2026", 0.7, 2.75, 12.0, 0.7,
    size=24, color=RGBColor(0xCC, 0xD8, 0xF0), align=PP_ALIGN.CENTER)
txb(sl, "Omri Segev  |  Advisors: Ofir Lindenbaum, Bracha Laufer-Goldshtein, Amir Globerson",
    0.7, 3.55, 12.0, 0.6, size=16,
    color=RGBColor(0xAA, 0xBB, 0xDD), align=PP_ALIGN.CENTER)
for lbl, yp in [
    ("Method: H(n) Spectral Features + Nadler Fusion", 4.5),
    ("Results: Math | GPQA | Factual QA | RAG | Agentic", 5.1),
    ("Next: Phase 11a complete + Phase 11b code/navigation pilots", 5.7),
]:
    accent_box(sl, lbl, 2.5, yp, 8.3, 0.48,
               bg=RGBColor(0x00, 0x50, 0x9E), fg=WHITE, size=14)

# -----------------------------------------------------------------------
# SLIDE 2: WHAT ARE WE DETECTING -- ENTROPY TRACES
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "What Are We Detecting?",
          "MATH-500 / Qwen-7B T=1.0  |  3 correct (blue) vs 3 hallucinated (orange)  |  EPR = DC (mean) only -- discards all shape")
img(sl, os.path.join(PLOTS, "fig1_individual_traces.png"), 0.3, 1.2, 12.7, 5.9)
accent_box(sl,
    "Incorrect answers show HIGHER, MORE VARIABLE entropy throughout the reasoning trace. "
    "Gray-box method: only token log-probabilities from one forward pass -- no attention maps, no extra sampling.",
    0.3, 6.85, 12.7, 0.5, bg=NAVY, fg=WHITE, size=13, bold=True)

# -----------------------------------------------------------------------
# SLIDE 3: PSD -- HALLUCINATION CHANGES SPECTRAL SIGNATURE
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "Hallucination Changes the Spectral Signature",
          "Average PSD (DC/EPR removed)  |  4 panels: MATH-500 (1.5B, 7B) and GPQA (Mistral-7B, Qwen-7B)  |  T=1.0 solid, T=1.5 dashed")
img(sl, os.path.join(PLOTS, "psd_comparison.png"), 0.3, 1.2, 12.7, 5.4)
txb(sl,
    "x-axis: normalised frequency in cycles per token  (0 = DC removed by construction; 0.5 = alternates every 2 tokens)\n"
    "Low band (0-0.1): one full oscillation per 10+ tokens -- slow structural patterns, reasoning stages\n"
    "High band (0.4-0.5): rapid 2-3 token bursts -- local uncertainty spikes, hallucination signature",
    0.3, 6.7, 9.5, 0.78, size=10.5, color=DGREY)
accent_box(sl, "Pattern is CONSISTENT across 4 model/dataset combos: incorrect = more high-freq energy",
           9.9, 6.7, 3.2, 0.78, bg=NAVY, fg=WHITE, size=11, bold=True)

# -----------------------------------------------------------------------
# SLIDE 4: FEATURE LIBRARY
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "16 Spectral Features -- Individual Performance on MATH-500",
          "Phase 5 bar shows 12 features; 4 more added in Phase 8: cusum_max, pe_mean, hurst, cusum_shift_idx  |  16 total")
img(sl, os.path.join(PLOTS, "fig3_feature_aucs.png"), 0.3, 1.2, 7.8, 5.6)
# Right column: metadata + Nadler subset + key formulas
txb(sl, "Dataset:", 8.35, 1.25, 4.85, 0.35, size=11, bold=True, color=NAVY)
txb(sl,
    "MATH-500 (500 problems) / Qwen2.5-Math-7B-Instruct / T=1.0\n"
    "Model accuracy 68.7%  =>  ~310 correct, ~190 wrong  (balanced)\n"
    "Why 7B? 1.5B at T=1.0 had too few errors for stable detection\n"
    "(advisor feedback from previous meeting on class imbalance)",
    8.35, 1.60, 4.85, 0.95, size=9.5, color=DGREY)
accent_box(sl,
    "Best Nadler subset (4 of 16):\ntrace_length + spectral_centroid + rpdi + sw_var_peak\n"
    "Selected by rho-filtered exhaustive search (|rho|<0.75)\n=> 90.0% AUROC",
    8.35, 2.65, 4.85, 1.1, bg=NAVY, fg=WHITE, size=11, bold=False)
txb(sl, "Key feature formulas:", 8.35, 3.88, 4.85, 0.35, size=11, bold=True, color=NAVY)
img(sl, FORMULA_PANEL, 8.35, 4.25, 4.85, 2.55)

# -----------------------------------------------------------------------
# SLIDE 5: FEATURE CORRELATION TOPOLOGY (META-ANALYSIS)
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "Feature Correlation Topology -- 7,001 Samples Across 5 Domains",
          "All 16 features (incl. Phase 8 additions)  |  meta-analysis  |  Spearman rho  |  why Nadler works: orthogonality to EPR cluster")
img(sl, os.path.join(PLOTS, "meta_analysis_cell06_out0.png"), 0.8, 1.25, 8.7, 6.0)
txb(sl, "What this shows:", 9.8, 1.5, 3.4, 0.4, size=13, bold=True, color=NAVY)
bullet_box(sl, [
    ("->", "epr / hurst / pe_mean are correlated\n   (all measure global entropy level)", DGREY),
    ("->", "spectral_entropy and stft_entropy correlated\n   (same signal, different window)", DGREY),
    ("->", "sw_var_peak + cusum_max are ORTHOGONAL\n   to EPR cluster => Nadler gains signal", GREEN),
    ("->", "pe_min isolated from all clusters\n   => confirmed noise, removed (rank 17/17)", RED),
], 9.8, 2.0, 3.4, 4.2, size=11.5)

# -----------------------------------------------------------------------
# SLIDE 6: FEATURE IMPORTANCE PER DOMAIN (META-ANALYSIS)
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "Feature Importance Varies By Domain -- Meta-Analysis (7,001 Samples)",
          "Random Forest Gini importance  |  5 domains  |  bottom-right panel: cross-domain verdict")
img(sl, FEAT_GRID, 0.2, 1.25, 12.9, 6.0)

# -----------------------------------------------------------------------
# SLIDE 7: METHOD APPLIED -- MATH EXAMPLE (NEW)
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "How It Works: Math Reasoning Domain",
          "H(n) covers the FULL generation trace (reasoning chain + answer)  |  ~500-2000 tokens")
img(sl, MATH_EX, 0.3, 1.2, 12.8, 5.85)
accent_box(sl,
    "Hallucinating model: regime shifts (cusum_max) and local variance spikes (sw_var_peak) appear mid-reasoning. "
    "Correct model: smooth low-entropy trajectory.",
    0.3, 7.1, 12.8, 0.33, bg=NAVY, fg=WHITE, size=11.5)

# -----------------------------------------------------------------------
# SLIDE 8: MATH RESULTS + COMPARISON TABLE
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "Math Reasoning: Results vs Competitors",
          "Long reasoning traces (500-2000 tokens) give FFT sufficient frequency resolution")
img(sl, os.path.join(PLOTS, "fig5_avg_trajectories.png"), 0.3, 1.25, 6.7, 4.3)
txb(sl, "Performance comparison:", 7.2, 1.3, 6.0, 0.4, size=13, bold=True, color=NAVY)
simple_table(sl,
    headers=["Method", "Dataset / Model", "AUROC", "Labels", "Access", "Passes"],
    rows=[
        ["Spectral Nadler (Ours)", "MATH-500 / Qwen-7B T=1.0", "90.0%", "None", "Gray-box", "1"],
        ["Spectral Nadler (Ours)", "GSM8K / Llama-3.1-8B",     "76.0%", "None", "Gray-box", "1"],
        ["LapEigvals (2502.17598)", "GSM8K / Llama-3.1-8B",    "72.0%", "None", "White-box\n(attn.maps)", "1"],
    ],
    l=7.2, t=1.8, w=6.0, h=3.1,
    col_widths=[1.7, 1.7, 0.7, 0.65, 0.7, 0.65],
    hdr_size=10, row_size=9)
accent_box(sl,
    "+4.0 pp over closest unsupervised competitor on same task+model\n"
    "LapEigvals is WHITE-BOX: needs full attention maps (all layers+heads) + supervised logistic probe. "
    "Our method: token log-probs only (gray-box) -- strictly less model access.\n"
    "T=1.5 gives 96.6% -- ceiling artifact (more mistakes => easier task); honest comparison is T=1.0.",
    7.2, 5.15, 6.0, 1.25,
    bg=RGBColor(0xFF, 0xF0, 0xCC), fg=RGBColor(0x80, 0x40, 0x00), size=10.5)

# -----------------------------------------------------------------------
# SLIDE 9: GPQA RESULTS + COMPARISON TABLE
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "GPQA Diamond: Scale Matters -- PhD-Level Science MCQ",
          "198 questions, 4-option MCQ  |  traces shorter (~100-400 tokens)  |  detection ceiling = model's own task accuracy")
img(sl, os.path.join(PLOTS, "B3_gpqa_trajectories.png"), 0.3, 1.25, 7.5, 3.5)
img(sl, GPQA_BAR, 0.3, 4.8, 4.7, 2.55)
txb(sl, "Performance comparison:", 8.1, 1.3, 5.1, 0.4, size=13, bold=True, color=NAVY)
simple_table(sl,
    headers=["Method", "Model", "AUROC", "Labels", "Access", "Passes"],
    rows=[
        ["Spectral Nadler (Ours)", "Qwen-72B AWQ", "69.0%", "None", "Gray-box", "1"],
        ["Spectral Nadler (Ours)", "7B models",    "65.4%", "None", "Gray-box", "1"],
        ["No published spectral\ncompetitor on GPQA", "--", "--", "--", "--", "--"],
    ],
    l=8.1, t=1.8, w=5.1, h=2.4,
    col_widths=[2.0, 0.9, 0.65, 0.6, 0.65, 0.6],
    hdr_size=10, row_size=9)
bullet_box(sl, [
    ("->", "7B models: ~30-40% task accuracy => too few correct examples for detector", DGREY),
    ("->", "72B AWQ raises accuracy to 40.4% => +3.6 pp AUROC", DGREY),
    ("->", "Bottleneck: model performance, not the detector", DGREY),
    ("->", "Gate not yet passed: need >=70%; direction: even stronger model or ensembling", RED),
], 8.1, 4.35, 5.1, 3.0, size=11)

# -----------------------------------------------------------------------
# SLIDE 10: NEGATIVE RESULT -- FACTUAL QA
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "Negative Result: Spectral Features Don't Transfer to Factual QA",
          "Falcon-3-10B  |  300 samples each  |  CoT extends traces to 200-500 tokens -- still fails  |  TriviaQA plot only")
img(sl, os.path.join(PLOTS, "direct_vs_cot_comparison.png"), 0.5, 1.25, 6.3, 4.9)
txb(sl, "Comparison:", 7.1, 1.3, 6.0, 0.4, size=13, bold=True, color=NAVY)
simple_table(sl,
    headers=["Method", "TriviaQA", "WebQ CoT", "Labels", "Access", "Passes"],
    rows=[
        ["Spectral Nadler (Ours, CoT)", "53.6%", "61.9%", "None", "Gray-box", "1"],
        ["EPR direct answer (no CoT)",  "79.1%", "71.8%", "None", "Gray-box", "1"],
    ],
    l=7.1, t=1.8, w=6.0, h=2.0,
    col_widths=[1.9, 0.8, 0.7, 0.65, 0.7, 0.65],
    hdr_size=10, row_size=9)
accent_box(sl,
    "WebQ direct answer: 0 correct samples in run => AUC undefined for direct-answer mode.\n"
    "Plot above shows TriviaQA only. WebQ CoT = 61.9% (CoT added enough tokens to compute spectral features).",
    7.1, 3.97, 6.0, 0.85,
    bg=RGBColor(0xFF, 0xF0, 0xCC), fg=RGBColor(0x80, 0x40, 0x00), size=10.5)
bullet_box(sl, [
    ("->", "Direct traces: 20-50 tokens => FFT has no frequency resolution", RED),
    ("->", "CoT traces: 200-500 tokens => no systematic spectral structure", DGREY),
    ("->", "CoT SMOOTHS entropy (confidence masking) => worse than direct EPR", DGREY),
    ("->", "Factual recall != reasoning: model retrieves, not generates structure", DGREY),
], 7.1, 4.95, 6.0, 1.7, size=11)
accent_box(sl,
    "This is a BOUNDARY CONDITION, not a failure.\n"
    "Spectral features detect generative uncertainty during REASONING -- not factual recall.",
    7.1, 6.1, 6.0, 1.15, bg=NAVY, fg=WHITE, size=12)

# -----------------------------------------------------------------------
# SLIDE 11: RAG METHOD / CITATION EXAMPLE
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "How It Works: RAG Domain -- Detecting Ungrounded Statements",
          "L-CiteEval  |  model prompted to cite retrieved passages [n]  |  H(n) sliced by citation boundaries")
img(sl, RAG_EXAMPLE, 0.3, 1.25, 12.8, 5.9)

# -----------------------------------------------------------------------
# SLIDE 12: RAG RESULTS + COMPARISON TABLE
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "RAG Results: 16 Cells (4 Models x 4 Datasets)",
          "Median AUROC ~72.8%  |  12/16 cells >=70%  |  Best: llama8b/HotpotQA = 87.7%")
img(sl, os.path.join(PLOTS, "A_headline_auc_heatmap.png"), 0.3, 1.25, 7.4, 4.7)
txb(sl, "vs competitors:", 7.95, 1.3, 5.2, 0.4, size=13, bold=True, color=NAVY)
simple_table(sl,
    headers=["Method", "Dataset", "AUROC", "Labels", "Access", "Passes"],
    rows=[
        ["Spectral Nadler\n(Ours, Llama-8B)", "HotpotQA",      "87.7%",  "None",     "Gray-box", "1"],
        ["Spectral Nadler\n(Ours, Qwen-7B)",  "2Wiki",          "80.5%",  "None",     "Gray-box", "1"],
        ["LOS-Net (2503.14043)\nDIFFERENT TASK", "Std HotpotQA\n(no citations)", "72.92%", "Required", "Gray-box", "1"],
    ],
    l=7.95, t=1.8, w=5.2, h=2.9,
    col_widths=[1.5, 1.1, 0.7, 0.7, 0.65, 0.65],
    hdr_size=9, row_size=9)
bullet_box(sl, [
    ("*", "Novel task: unsupervised citation grounding on L-CiteEval -- no direct published competitor", GREEN),
    ("->", "LOS-Net uses standard HotpotQA (raw QA, no citation markers) -- different task", ORANGE),
    ("->", "Closest available benchmark: LOS-Net = 72.92%; we beat it by +14.8 pp (unsupervised!)", GREEN),
    ("->", "Median 16 cells: ~72.8% | 12/16 >= 70%", DGREY),
], 7.95, 4.9, 5.2, 1.7, size=11)

# -----------------------------------------------------------------------
# SLIDE 13: RAG SANITY CHECK
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "RAG Sanity: Signal Is Spectral Shape, Not Response Length",
          "trace_length alone gives 50.8% (chance) -- the AUC comes from spectral structure, not verbosity")
img(sl, os.path.join(PLOTS, "E_length_controlled_bars.png"), 0.4, 1.25, 12.5, 5.9)

# -----------------------------------------------------------------------
# SLIDE 14: RAG SCORE DISTRIBUTIONS
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "RAG: Score Distributions -- Grounded vs Ungrounded Statements",
          "Nadler-fused scores per (model, dataset) cell  |  AUC annotated per panel  |  separation is visible")
img(sl, os.path.join(PLOTS, "D_fusion_distributions_grid.png"), 0.4, 1.25, 12.5, 5.9)

# -----------------------------------------------------------------------
# SLIDE 15: HOW IT WORKS -- AGENTIC (NEW)
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "How It Works: Agentic ReAct Loops",
          "Phase 11a  |  Task: multi-hop QA via ReAct  |  H(n) from Thought tokens only  |  Phi_min = weakest-link aggregation")
# Task summary box (left column, compact)
txb(sl, "Task (Phase 11a):", 0.3, 1.25, 6.0, 0.35, size=11, bold=True, color=NAVY)
txb(sl,
    "Input: question + Wikipedia passage dictionary\n"
    "Model: ReAct loop, up to 3 steps (Thought => Action: search/finish => Observation)\n"
    "Label: trajectory_correct = final answer matches gold answer string",
    0.3, 1.62, 6.2, 0.75, size=10, color=DGREY)
# RAG vs Agentic comparison
accent_box(sl,
    "vs RAG (Phase 10):  RAG = 1 forward pass, model generates full cited response, H(n) sliced by citation markers.\n"
    "Agentic = 3 separate passes; H(n) extracted from each Thought step (50-150 tokens); then aggregated via Phi.",
    0.3, 2.47, 6.2, 0.85, bg=RGBColor(0xE8, 0xF0, 0xFE), fg=DGREY, size=10)
img(sl, AGENTIC_EX, 0.3, 3.42, 12.8, 3.72)
accent_box(sl,
    "Phi_min = min(score_step1, score_step2, score_step3) -- weakest link.  "
    "Phi_avg = mean across steps.  Phi_last = final Thought only.\n"
    "Phi_min outperforms: a reasoning chain is only as reliable as its most uncertain step.",
    0.3, 7.2, 12.8, 0.25, bg=NAVY, fg=WHITE, size=11)

# -----------------------------------------------------------------------
# SLIDE 16: AGENTIC RESULTS + COMPARISON TABLE
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "Agentic: Early Signal vs AUQ Baseline (Zhang et al. 2026)",
          "4 models x 2 datasets x 3 aggregations  |  run not yet complete (Mistral-24B + Qwen-72B pending)")
# Phi definitions box (top left, compact)
accent_box(sl,
    "Phi_min = min(score_step1, …, score_stepK)   [weakest link]\n"
    "Phi_avg = mean(score_step1, …, score_stepK)   [average reliability]\n"
    "Phi_last = score_stepK only   [final Thought quality]",
    0.3, 1.25, 6.1, 1.0,
    bg=RGBColor(0xE8, 0xF0, 0xFE), fg=DGREY, size=11)
img(sl, AGENTIC_BAR, 0.3, 2.35, 5.8, 5.0)
txb(sl, "Comparison:", 6.6, 1.3, 6.5, 0.4, size=13, bold=True, color=NAVY)
simple_table(sl,
    headers=["Method", "Setting", "AUROC", "Labels", "Access", "Passes"],
    rows=[
        ["Spectral Nadler Phi_min\n(Ours, mid-run)",  "2Wiki / DeepSeek-7B", "85.0%*", "None", "Gray-box",   "1/step"],
        ["AUQ verbalized\n(Zhang et al. 2026)",        "ALFWorld / Phi_min",  "79.1%",  "None", "White-box",  "1/step"],
    ],
    l=6.6, t=1.8, w=6.5, h=2.3,
    col_widths=[1.9, 1.4, 0.75, 0.65, 0.8, 0.75],
    hdr_size=10, row_size=9)
bullet_box(sl, [
    ("->", "4 models: Qwen-7B, DeepSeek-R1-7B, Mistral-24B, Qwen-72B", DGREY),
    ("->", "Datasets: HotpotQA + 2WikiMultiHopQA (N=200 each)", DGREY),
    ("->", "AUQ = asks model 'how confident are you?' in natural language", ORANGE),
    ("->", "RLHF-aligned models often claim high confidence even when wrong", ORANGE),
    ("->", "Our method: no verbalization; reads spectral shape only => gray-box", GREEN),
], 6.6, 4.3, 6.5, 2.1, size=11)
accent_box(sl,
    "Run not yet complete: Mistral-24B + Qwen-72B cells pending.\n"
    "Mid-run signal is DeepSeek-R1-7B only. Official results after full run.",
    6.6, 6.55, 6.5, 0.7,
    bg=RGBColor(0xFF, 0xF0, 0xCC), fg=RGBColor(0x80, 0x40, 0x00), size=11)

# -----------------------------------------------------------------------
# SLIDE 17: WHAT'S NEXT
# -----------------------------------------------------------------------
sl = add_slide()
title_bar(sl, "What's Next",
          "Phase 11a completion + Phase 11b new domains + thesis end-game")

# Four panels: Phase 11a, Phase 11b Pilot A, Phase 11b Pilot B, bottom end-game
panel_data = [
    ("Phase 11a: Complete ReAct Multi-hop QA", BLUE, [
        "Run Mistral-24B (normal runtime)",
        "Run Qwen-72B (fresh runtime + gptqmodel stub)",
        "Full analysis: 8 cells (4 models x 2 datasets)",
        "Official Phi_min / Phi_avg / Phi_last AUC table",
        "",
        "Competitor: AUQ (Zhang et al. 2026)",
        "  White-box verbalized confidence = 79.1%",
        "  Our mid-run gray-box: 85.0% (+5.9 pp)",
        "",
        "Target: beat AUQ by >=5 pp, unsupervised",
    ]),
    ("Phase 11b Pilot A: HumanEval (Code)", PURPLE, [
        "Domain: code generation + execution",
        "N=20 problems, qwen25_7b, 3 attempts each",
        "Label: unit test pass/fail (objective!)",
        "H(n) from code generation tokens",
        "",
        "Competitor: DSDE (2026)",
        "  Execution-based disagreement",
        "  AUROC: 0.82-0.84",
        "",
        "GO/NO-GO: >=5 solved, traces non-degenerate",
    ]),
    ("Phase 11b Pilot B: ALFWorld (Navigation)", ORANGE, [
        "Domain: embodied text navigation",
        "N=5 tasks, pick_and_place type",
        "Model: qwen25_7b, MAX_STEPS=20",
        "Label: task_success (binary)",
        "H(n) from Thought tokens per step",
        "",
        "Competitor: AUQ (Zhang et al. 2026)",
        "  Same paper as Phase 11a SOTA",
        "  AUROC: 0.791 on ALFWorld",
        "",
        "GO/NO-GO: alfworld imports stable",
    ]),
]

for ci, (panel_title, color, items) in enumerate(panel_data):
    lx = 0.3 + ci * 4.35
    accent_box(sl, panel_title, lx, 1.3, 4.1, 0.6, bg=color, fg=WHITE, size=12, bold=True)
    tb = sl.shapes.add_textbox(Inches(lx), Inches(2.05), Inches(4.1), Inches(4.2))
    tf = tb.text_frame; tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.add_paragraph() if j > 0 else tf.paragraphs[0]
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = ("    " + item) if item.startswith("  ") else ("- " + item if item else "")
        run.font.size = Pt(10.5)
        run.font.color.rgb = RED if 'Competitor' in item or 'AUROC' in item or 'DSDE' in item or 'AUQ' in item else DGREY

# Bottom: thesis end-game
accent_box(sl, "Thesis end-game (after Phase 11):", 0.3, 6.35, 12.75, 0.4,
           bg=NAVY, fg=WHITE, size=12, bold=True)
for ci, (title, body) in enumerate([
    ("LTT Conformal Calibration (Bracha)",
     "~50 lines | data exists | converts AUROC -> deployment guarantee: recall>=90% at 95% confidence"),
    ("Manifold Analysis (Ofir)",
     "Do entropy trajectories lie on a manifold? Hallucination = escape. LOCA + IMM."),
    ("Update Research_Directions + thesis write-up",
     "4 arenas confirmed: math | GPQA | RAG | agentic. 2 new pilots: code + navigation."),
]):
    lx = 0.3 + ci * 4.35
    txb(sl, title, lx, 6.82, 4.1, 0.3, size=11, bold=True, color=NAVY)
    txb(sl, body,  lx, 7.14, 4.1, 0.32, size=9.5, color=DGREY)

# =======================================================================
prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"  {len(prs.slides)} slides")
